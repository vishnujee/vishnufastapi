import os
import tempfile
import logging
import boto3
import hashlib
import io
import zipfile
import subprocess
import psutil
import fitz
from PyPDF2 import PdfMerger, PdfReader, PdfWriter
from PIL import Image
import pdfplumber
from pdf2docx import Converter
import pandas as pd
import time
import gc
import platform
from pptx import Presentation 
from pdf2image import convert_from_bytes

from typing import Dict,  Optional

from reportlab.lib.utils import ImageReader
from reportlab.pdfgen import canvas
import logging
from docx import Document

# Configure logging
logger = logging.getLogger(__name__)

from reportlab.pdfgen import canvas


from rembg import remove

logger = logging.getLogger(__name__)

from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import math

# Configure logging
os.makedirs("logs", exist_ok=True)
logging.basicConfig(
    level=logging.INFO,
    filename=os.path.join("logs", "pdfoperations.log"),
    format="%(asctime)s - %(levelname)s - %(message)s"
)
logger = logging.getLogger(__name__)

# S3 and local mode
BUCKET_NAME = os.getenv("BUCKET_NAME")
AWS_ACCESS_KEY = os.getenv("AWS_ACCESS_KEY_ID")
AWS_SECRET_KEY = os.getenv("AWS_SECRET_ACCESS_KEY")



USE_S3 = all([BUCKET_NAME, AWS_ACCESS_KEY, AWS_SECRET_KEY])

if USE_S3:
    # Split bucket name and folder prefix (if provided)
    if "/" in BUCKET_NAME:
        BUCKET_NAME, S3_PREFIX = BUCKET_NAME.split("/", 1)
    else:
        S3_PREFIX = ""

    # Initialize S3 client
    s3_client = boto3.client(
        "s3",
        aws_access_key_id=AWS_ACCESS_KEY,
        aws_secret_access_key=AWS_SECRET_KEY,
    )
else:
    print("AWS credentials missing. Falling back to local storage.")
    os.makedirs("input_pdfs", exist_ok=True)
    os.makedirs("output_pdfs", exist_ok=True)


# s3_client = boto3.client(
#         "s3",
#         aws_access_key_id=AWS_ACCESS_KEY,
#         aws_secret_access_key=AWS_SECRET_KEY,
#     )



def upload_to_s3(file_content, filename):
    """Upload file content to S3 or save locally."""

    s3_key = f"temp_uploads/{hashlib.md5(file_content).hexdigest()}_{filename}"
    s3_client.put_object(Bucket=BUCKET_NAME, Key=s3_key, Body=file_content)
    logger.info(f"Uploaded to S3: {s3_key}")
    return s3_key

def download_from_s3(s3_key):
    """Download file content from S3 or read locally."""

    try:
        response = s3_client.get_object(Bucket=BUCKET_NAME, Key=s3_key)
        return response["Body"].read()
    except Exception as e:
        logger.error(f"Failed to download S3 file {s3_key}: {e}")
        raise

def cleanup_s3_file(s3_key):
    """Delete file from S3 or locally."""

    try:
        s3_client.delete_object(Bucket=BUCKET_NAME, Key=s3_key)
        logger.info(f"Deleted S3 file: {s3_key}")
    except Exception as e:
        logger.warning(f"Failed to delete S3 file {s3_key}: {e}")

def get_memory_info():
    """Return system-wide memory usage info in MB."""
    try:
        total_memory = psutil.virtual_memory().total / (1024 * 1024)
        used_memory = psutil.virtual_memory().used / (1024 * 1024)
        available_memory = psutil.virtual_memory().available / (1024 * 1024)
        return used_memory, available_memory, total_memory
    except Exception as e:
        logger.error(f"Failed to get memory info: {e}")
        return 0, 0, 0

def get_process_memory(pid):
    """Return memory usage of the current Python process in MB."""
    try:
        process = psutil.Process(pid)
        memory_info = process.memory_info()
        return memory_info.rss / (1024 * 1024)
    except (psutil.NoSuchProcess, psutil.AccessDenied, OSError):
        return None

def get_subprocess_memory(pid):
    """Get memory usage of a subprocess in MB."""
    try:
        process = psutil.Process(pid)
        memory_info = process.memory_info()
        return memory_info.rss / (1024 * 1024)
    except (psutil.NoSuchProcess, psutil.AccessDenied):
        return None

def merge_pdfs_pypdf2(file_contents):
    """Merge PDFs using PyPDF2 with disk-based operations."""
    merger = PdfMerger()
    temp_files = []
    try:
        # Write each PDF to temp file
        for content in file_contents:
            with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
                tmp.write(content)
                temp_files.append(tmp.name)
        
        # Append to merger (reads sequentially from disk)
        for tmp_file in temp_files:
            merger.append(tmp_file)
        
        # Write merged PDF to temp file first
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp_output:
            merger.write(tmp_output.name)
            with open(tmp_output.name, "rb") as f:
                return f.read()
    except Exception as e:
        logger.error(f"PyPDF2 merge error: {e}")
        return None
    finally:
        merger.close()
        for tmp_file in temp_files:
            try:
                os.unlink(tmp_file)
            except:
                pass
        gc.collect()
def merge_pdfs_ghostscript(file_contents, output_path):
    """Merge PDFs using Ghostscript."""
    input_files = []
    try:
        total_size_mb = sum(len(content) for content in file_contents) / (1024 * 1024)
        if total_size_mb > 50:
            raise Exception(f"Total file size {total_size_mb:.2f}MB exceeds 50MB limit")
        
        for content in file_contents:
            with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
                tmp.write(content)
                input_files.append(tmp.name)
        
        # Determine Ghostscript binary based on OS
        gs_binary = "gswin64c" if platform.system() == "Windows" else "gs"
        
        gs_command = [
            gs_binary,
            "-sDEVICE=pdfwrite",
            "-dCompatibilityLevel=1.4",
            "-dPDFSETTINGS=/default",
            "-dNOPAUSE",
            "-dBATCH",
            f"-sOutputFile={output_path}"
        ] + input_files

        logger.info(f"Running Ghostscript command: {' '.join(gs_command)}")
        process = subprocess.run(gs_command, capture_output=True, text=True)
        if process.returncode != 0:
            logger.error(f"Ghostscript error: {process.stderr}")
            return None

        with open(output_path, "rb") as f:
            return f.read()
    except Exception as e:
        logger.error(f"Ghostscript merge error: {e}")
        return None
    finally:
        for f in input_files:
            if os.path.exists(f):
                os.unlink(f)
        if os.path.exists(output_path):
            os.unlink(output_path)
        gc.collect()

def safe_compress_pdf(pdf_bytes, dpi, quality):
    """Compress PDF using Ghostscript with memory monitoring."""
    input_file = None
    output_file = None
    gs_binary = "gswin64c" if platform.system() == "Windows" else "gs"
    try:
        # Verify Ghostscript is installed
        result = subprocess.run([gs_binary, "--version"], capture_output=True, text=True)
        if result.returncode != 0:
            raise FileNotFoundError(f"Ghostscript ('{gs_binary}') is not installed or not found in PATH.")

        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp_input:
            tmp_input.write(pdf_bytes)
            tmp_input.flush()
            input_file = tmp_input.name

        output_file = tempfile.NamedTemporaryFile(delete=False, suffix=".pdf").name

        pdf_settings = "/screen" if dpi <= 72 else "/ebook"
        gs_command = [
            gs_binary,
            "-sDEVICE=pdfwrite",
            "-dCompatibilityLevel=1.4",
            f"-dPDFSETTINGS={pdf_settings}",
            "-dColorImageDownsampleType=/Bicubic",
            f"-dColorImageResolution={dpi}",
            "-dGrayImageDownsampleType=/Bicubic",
            f"-dGrayImageResolution={dpi}",
            "-dMonoImageDownsampleType=/Subsample",
            f"-dMonoImageResolution={dpi}",
            "-dColorImageFilter=/DCTEncode",
            "-dGrayImageFilter=/DCTEncode",
            "-dNOPAUSE",
            "-dBATCH",
            "-sOutputFile=" + output_file,
            input_file
        ]

        peak_gs_memory = 0.0
        process = subprocess.Popen(
            gs_command,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True
        )

        while process.poll() is None:
            used_mem, avail_mem, total_mem = get_memory_info()
            mem_percent = (used_mem / total_mem) * 100
            gs_memory = get_subprocess_memory(process.pid)
            if gs_memory is not None:
                peak_gs_memory = max(peak_gs_memory, gs_memory)
            if mem_percent > 80:
                process.terminate()
                try:
                    process.wait(timeout=2)
                except subprocess.TimeoutExpired:
                    process.kill()
                raise Exception("High memory usage detected")
            time.sleep(0.5)

        stdout, stderr = process.communicate()
        if process.returncode != 0:
            logger.error(f"Ghostscript error: {stderr}")
            return None

        with open(output_file, "rb") as f:
            compressed_pdf = f.read()

        return compressed_pdf
    except FileNotFoundError as e:
        logger.error(f"Compression error: {e}")
        return None
    except Exception as e:
        logger.error(f"Compression error: {e}")
        return None
    finally:
        for path in [input_file, output_file]:
            if path and os.path.exists(path):
                try:
                    os.unlink(path)
                except:
                    pass
        gc.collect()

def encrypt_pdf(pdf_bytes, password):
    """Encrypt PDF with a password using PyMuPDF."""
    try:
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        used_mem, avail_mem, total_mem = get_memory_info()
        mem_percent = (used_mem / total_mem) * 100
        current_pid = os.getpid()
        process_memory = get_process_memory(current_pid)
        if mem_percent > 70:
            raise Exception("High memory usage detected")
        
        output_buffer = io.BytesIO()
        doc.save(
            output_buffer,
            encryption=fitz.PDF_ENCRYPT_AES_256,
            owner_pw=password,
            user_pw=password,
            permissions=fitz.PDF_PERM_PRINT
        )
        doc.close()
        output_buffer.seek(0)
        
        encrypted_bytes = output_buffer.getvalue()
        test_doc = fitz.open(stream=encrypted_bytes, filetype="pdf")
        is_encrypted = test_doc.is_encrypted
        test_doc.close()
        
        if not is_encrypted:
            raise Exception("Encryption failed: Output PDF is not encrypted")
        
        return encrypted_bytes
    except Exception as e:
        logger.error(f"Encryption error: {e}")
        return None
    finally:
        gc.collect()

def convert_pdf_to_images(pdf_bytes):
    """Convert PDF pages to images and return a ZIP file."""
    try:
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        num_pages = len(doc)
        if num_pages == 0:
            raise Exception("The PDF is empty")
        
        used_mem, avail_mem, total_mem = get_memory_info()
        mem_percent = (used_mem / total_mem) * 100
        current_pid = os.getpid()
        process_memory = get_process_memory(current_pid)
        if mem_percent > 80:
            raise Exception("High memory usage detected")

        zip_buffer = io.BytesIO()
        with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zip_file:
            for page_num in range(num_pages):
                page = doc[page_num]
                pix = page.get_pixmap(dpi=150)
                img_bytes = pix.tobytes("png")
                zip_file.writestr(f"page_{page_num + 1}.png", img_bytes)
                
                used_mem, avail_mem, total_mem = get_memory_info()
                mem_percent = (used_mem / total_mem) * 100
                if mem_percent > 80:
                    raise Exception(f"High memory usage during page {page_num + 1}")
        
        doc.close()
        zip_buffer.seek(0)
        return zip_buffer.getvalue()
    except Exception as e:
        logger.error(f"PDF to images error: {e}")
        return None
    finally:
        gc.collect()

def split_pdf(pdf_bytes):
    """Split PDF into individual pages and return a ZIP file."""
    try:
        reader = PdfReader(io.BytesIO(pdf_bytes))
        total_pages = len(reader.pages)
        zip_buffer = io.BytesIO()
        with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zipf:
            for i, page in enumerate(reader.pages):
                writer = PdfWriter()
                writer.add_page(page)
                page_bytes = io.BytesIO()
                writer.write(page_bytes)
                page_bytes.seek(0)
                zipf.writestr(f"page_{i+1}.pdf", page_bytes.read())
        zip_buffer.seek(0)
        return zip_buffer.getvalue(), total_pages
    except Exception as e:
        logger.error(f"Split error: {e}")
        return None, 0
    finally:
        gc.collect()

def delete_pdf_pages(pdf_bytes, pages_to_delete):
    """Delete specified pages from PDF."""
    try:
        pdf_reader = PdfReader(io.BytesIO(pdf_bytes))
        total_pages = len(pdf_reader.pages)
        if not pages_to_delete or max(pages_to_delete) > total_pages:
            raise Exception(f"Invalid pages to delete. PDF has {total_pages} pages")
        
        pdf_writer = PdfWriter()
        for page_num in range(1, total_pages + 1):
            if page_num not in pages_to_delete:
                pdf_writer.add_page(pdf_reader.pages[page_num - 1])
        
        output_buffer = io.BytesIO()
        pdf_writer.write(output_buffer)
        output_buffer.seek(0)
        return output_buffer.getvalue()
    except Exception as e:
        logger.error(f"Delete pages error: {e}")
        return None
    finally:
        gc.collect()

def convert_pdf_to_word(pdf_bytes):
    """Convert PDF to Word using pdf2docx."""
    input_file = None
    output_file = None
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp_input:
            tmp_input.write(pdf_bytes)
            input_file = tmp_input.name
        
        output_file = tempfile.NamedTemporaryFile(delete=False, suffix=".docx").name
        cv = Converter(input_file)
        cv.convert(output_file)
        cv.close()
        
        with open(output_file, "rb") as f:
            return f.read()
    except Exception as e:
        logger.error(f"PDF to Word error: {e}")
        return None
    finally:
        for path in [input_file, output_file]:
            if path and os.path.exists(path):
                try:
                    os.unlink(path)
                except:
                    pass
        gc.collect()


def convert_pdf_to_ppt(pdf_bytes):
    """Convert PDF to PowerPoint with proper aspect ratio handling"""
    try:
        # Convert PDF to images
        images = convert_from_bytes(
            pdf_bytes,
            dpi=200,
            fmt='jpeg',
            jpegopt={'quality': 90},
            thread_count=4
        )
        
        if not images:
            raise ValueError("No pages converted from PDF")

        prs = Presentation()
        
        # Set to standard 4:3 aspect ratio (more square than 16:9)
        prs.slide_width = Inches(7)       # Was 10
        prs.slide_height = Inches(5.25)     # Reduced from 7.5 inches# 7.5 inches tall (4:3 ratio)
        
        # Available content area with margins
        content_width = Inches(6.5)       # Was 9 
        content_height = Inches(4.75)     # Was 6.75     # Reduced from 6.75 inches       # 0.375" top + bottom margins

        for img in images:
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
            
            # Get original image dimensions
            img_width, img_height = img.size
            
            # Calculate scaling to fit content area while maintaining aspect ratio
            width_ratio = content_width / img_width
            height_ratio = content_height / img_height
            scale = min(width_ratio, height_ratio)
            
            # Apply scaling
            scaled_width = img_width * scale
            scaled_height = img_height * scale
            
            # Center the image on slide
            left = (prs.slide_width - scaled_width) / 2
            top = (prs.slide_height - scaled_height) / 2
            
            # Convert image to bytes
            with io.BytesIO() as output:
                img.save(output, format='JPEG', quality=90)
                output.seek(0)
                
                # Add to slide
                slide.shapes.add_picture(
                    output,
                    left,
                    top,
                    width=scaled_width,
                    height=scaled_height
                )

        # Save presentation
        output = io.BytesIO()
        prs.save(output)
        return output.getvalue()

    except Exception as e:
        logger.error(f"PPT conversion failed: {str(e)}")
        return None



def convert_pdf_to_editable_ppt(pdf_bytes):
    """Convert PDF to properly formatted PowerPoint with error handling"""
    try:
        from pptx.util import Pt, Inches
        from pptx.enum.text import PP_ALIGN, PP_PARAGRAPH_ALIGNMENT
        from pptx.dml.color import RGBColor
        import re
        import tempfile
        import io
        import os
        from docx import Document
        from pdf2docx import Converter
        from pptx import Presentation
        import logging

        # Set up logging
        logging.basicConfig(level=logging.INFO)
        logger = logging.getLogger(__name__)

        # Create temp files
        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as pdf_temp:
            pdf_temp.write(pdf_bytes)
            pdf_path = pdf_temp.name
        
        docx_path = tempfile.mktemp(suffix=".docx")
        
        # Convert PDF to Word
        cv = Converter(pdf_path)
        cv.convert(docx_path, 
                  layout_analysis=True,
                  keep_blank_lines=False,
                  ignore_page_error=True)
        cv.close()
        
        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        # Style definitions
        def safe_pt(value, default=12):
            return Pt(value) if value is not None else Pt(default)
            
        title_style = {
            'size': safe_pt(36),
            'bold': True,
            'color': RGBColor(0, 32, 96)
        }
        content_style = {
            'size': safe_pt(18),  # Further reduced size
            'color': RGBColor(64, 64, 64),
            'alignment': PP_ALIGN.LEFT,
            'space_after': safe_pt(8),  # Reduced spacing
            'line_spacing': 1.2,  # Increased line spacing
            'word_wrap': True,
            'width': Inches(11.5),  # Text box width
            'height': Inches(5)     # Text box height
        }

        # Process document
        doc = Document(docx_path)
        current_slide_content = []
        slide_count = 0
        
        def add_safe_slide(content):
            """Create slide with proper text wrapping and boundaries"""
            nonlocal slide_count
            try:
                if not content.strip():
                    return
                    
                slide_count += 1
                slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
                
                # Add title
                title = slide.shapes.add_textbox(
                    left=Inches(0.5), top=Inches(0.25),
                    width=Inches(12), height=Inches(0.75))  # Smaller title box
                title.text = f"Slide {slide_count}"
                title.text_frame.word_wrap = True
                
                # Format title
                for paragraph in title.text_frame.paragraphs:
                    paragraph.font.size = title_style['size']
                    paragraph.font.bold = title_style['bold']
                    paragraph.font.color.rgb = title_style['color']
                    paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                
                # Add content with strict boundaries
                body = slide.shapes.add_textbox(
                    left=Inches(0.75), top=Inches(1.25),  # Adjusted position
                    width=content_style['width'], 
                    height=content_style['height'])
                
                # Clean content and handle long lines
                cleaned_content = []
                for line in content.split('\n'):
                    # Split long lines at approximately 100 characters
                    if len(line) > 100:
                        chunks = [line[i:i+100] for i in range(0, len(line), 100)]
                        cleaned_content.extend(chunks)
                    else:
                        cleaned_content.append(line)
                
                body.text = '\n'.join(cleaned_content)
                body.text_frame.word_wrap = True
                body.text_frame.auto_size = None  # Disable auto-size to enforce boundaries
                
                # Format content
                for paragraph in body.text_frame.paragraphs:
                    paragraph.font.size = content_style['size']
                    paragraph.font.color.rgb = content_style['color']
                    paragraph.alignment = content_style['alignment']
                    paragraph.space_after = content_style['space_after']
                    paragraph.line_spacing = content_style['line_spacing']
                    
                    # Remove any remaining special characters
                    paragraph.text = re.sub(r'[•\u2022\u25CF\uFEFF]', '', paragraph.text)
                    
            except Exception as e:
                logger.error(f"Slide creation error: {str(e)}")

        # Paragraph processing with line length control
        for para in doc.paragraphs:
            try:
                text = re.sub(r'[•\u2022\u25CF\uFEFF]', '', para.text).strip()
                if not text:
                    continue
                    
                # Split long paragraphs into chunks
                if len(text) > 300:
                    chunks = [text[i:i+300] for i in range(0, len(text), 300)]
                    current_slide_content.extend(chunks)
                else:
                    current_slide_content.append(text)
                
                # Check if we should create a new slide
                if len('\n'.join(current_slide_content)) > 300:
                    add_safe_slide('\n'.join(current_slide_content))
                    current_slide_content = []
                
            except Exception as e:
                logger.error(f"Paragraph processing error: {str(e)}")
                continue
        
        # Add final slide
        if current_slide_content:
            add_safe_slide('\n'.join(current_slide_content))
        
        # Validate presentation
        if slide_count == 0:
            raise ValueError("No valid slides created")
            
        # Save with checks
        output = io.BytesIO()
        prs.save(output)
        if output.getbuffer().nbytes < 1024:
            raise ValueError("Presentation too small - likely conversion failed")
            
        return output.getvalue()
        
    except Exception as e:
        logger.error(f"PPT conversion failed: {str(e)}")
        return None
    finally:
        # Cleanup
        for path in [pdf_path, docx_path]:
            try:
                if path and os.path.exists(path):
                    os.unlink(path)
            except:
                pass


def convert_pdf_to_ppt(pdf_bytes):
    """Convert PDF to PowerPoint with tight width and proper image centering"""
    try:
        # Convert PDF to images with error handling
        images = []
        try:
            images = convert_from_bytes(
                pdf_bytes,
                dpi=200,
                fmt='jpeg',
                jpegopt={'quality': 90},
                thread_count=4
            )
            if not images:
                raise ValueError("No pages converted from PDF")
        except Exception as e:
            raise ValueError(f"PDF to image conversion failed: {str(e)}")

        prs = Presentation()
        
        # Set ultra-compact slide dimensions (adjusted 16:9 ratio)
        slide_width = Inches(4.0)   # Tight width (~10.16 cm)
        slide_height = Inches(2.25) # Maintains 16:9 aspect ratio
        prs.slide_width = slide_width
        prs.slide_height = slide_height

        for img in images:
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

            with io.BytesIO() as output:
                img.save(output, format='JPEG', quality=90)
                output.seek(0)
                
                # Calculate with minimal padding
                img_ratio = img.width / img.height
                
                # Use 90% of slide width for image (5% margin each side)
                target_width = slide_width * 0.9
                target_height = target_width / img_ratio
                
                # If too tall, scale down to fit height instead
                if target_height > slide_height * 0.9:
                    target_height = slide_height * 0.9
                    target_width = target_height * img_ratio
                
                # Center the image
                left = (slide_width - target_width) / 2
                top = (slide_height - target_height) / 2

                slide.shapes.add_picture(
                    output,
                    left,
                    top,
                    width=target_width,
                    height=target_height
                )

        output = io.BytesIO()
        prs.save(output)
        return output.getvalue()

    except Exception as e:
        logger.error(f"PPT conversion failed: {str(e)}")
        return None


def convert_pdf_to_excel(pdf_bytes):
    """Convert PDF to Excel using pdfplumber."""
    try:
        with pdfplumber.open(io.BytesIO(pdf_bytes)) as pdf:
            all_tables = []
            for page in pdf.pages:
                tables = page.extract_tables()
                for table in tables:
                    all_tables.append(pd.DataFrame(table[1:], columns=table[0]))
        
        if not all_tables:
            raise Exception("No tables found in PDF")
        
        df = pd.concat(all_tables, ignore_index=True) if len(all_tables) > 1 else all_tables[0]
        output_buffer = io.BytesIO()
        df.to_excel(output_buffer, index=False, engine="openpyxl")
        output_buffer.seek(0)
        return output_buffer.getvalue()
    except Exception as e:
        logger.error(f"PDF to Excel error: {e}")
        return None
    finally:
        gc.collect()

def convert_image_to_pdf(image_bytes, page_size="A4", orientation="Portrait"):
    """Convert a single image to a one-page PDF."""
    try:
        doc = fitz.open()
        page_sizes = {"A4": (595, 842), "Letter": (612, 792)}
        page_width, page_height = page_sizes.get(page_size, page_sizes["A4"])
        if orientation == "Landscape":
            page_width, page_height = page_height, page_width
        
        used_mem, avail_mem, total_mem = get_memory_info()
        mem_percent = (used_mem / total_mem) * 100
        if mem_percent > 80:
            raise Exception("High memory usage detected")
        
        img = Image.open(io.BytesIO(image_bytes))
        if img.format not in ["PNG", "JPEG"]:
            raise Exception("Only PNG and JPEG are supported")
        img_width, img_height = img.size
        
        margin = 10
        usable_width = page_width - 2 * margin
        usable_height = page_height - 2 * margin
        scale = min(usable_width / img_width, usable_height / img_height)
        new_width = img_width * scale
        new_height = img_height * scale
        
        page = doc.new_page(width=page_width, height=page_height)
        x0 = (page_width - new_width) / 2
        y0 = (page_height - new_height) / 2
        rect = fitz.Rect(x0, y0, x0 + new_width, y0 + new_height)
        
        img_buffer = io.BytesIO()
        img.save(img_buffer, format="PNG")
        page.insert_image(rect, stream=img_buffer.getvalue())
        img.close()
        
        pdf_bytes = io.BytesIO()
        doc.save(pdf_bytes)
        doc.close()
        pdf_bytes.seek(0)
        return pdf_bytes.getvalue()
    except Exception as e:
        logger.error(f"Image to PDF error: {e}")
        return None
    finally:
        gc.collect()

def remove_pdf_password(pdf_bytes, password):
    """Remove password from PDF using PyMuPDF."""
    try:
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        if not doc.is_encrypted:
            return pdf_bytes
        if not doc.authenticate(password):
            raise Exception("Incorrect password")
        
        used_mem, avail_mem, total_mem = get_memory_info()
        mem_percent = (used_mem / total_mem) * 100
        if mem_percent > 80:
            raise Exception("High memory usage detected")
        
        output_buffer = io.BytesIO()
        doc.save(output_buffer, encryption=0)
        doc.close()
        output_buffer.seek(0)
        
        decrypted_bytes = output_buffer.getvalue()
        test_doc = fitz.open(stream=decrypted_bytes, filetype="pdf")
        is_encrypted = test_doc.is_encrypted
        test_doc.close()
        
        if is_encrypted:
            raise Exception("Failed to remove password")
        
        return decrypted_bytes
    except Exception as e:
        logger.error(f"PDF password removal error: {e}")
        return None
    finally:
        gc.collect()




def reorder_pdf_pages(pdf_bytes, page_order):
    try:
        pdf_reader = PdfReader(io.BytesIO(pdf_bytes))
        total_pages = len(pdf_reader.pages)
        
        if not page_order or len(page_order) != total_pages:
            raise Exception(f"Invalid page order. Must specify {total_pages} unique pages.")
        if sorted(page_order) != list(range(1, total_pages + 1)):
            raise Exception("Page order must include all pages exactly once.")
        
        pdf_writer = PdfWriter()
        for page_num in page_order:
            pdf_writer.add_page(pdf_reader.pages[page_num - 1])
        
        output_buffer = io.BytesIO()
        pdf_writer.write(output_buffer)
        output_buffer.seek(0)
        return output_buffer.getvalue()
    except Exception as e:
        logger.error(f"Page reordering error: {e}")
        return None
    finally:
        gc.collect()

def add_page_numbers(pdf_bytes, position="bottom", alignment="center", format="page_x"):
    try:
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        total_pages = len(doc)
        if position not in ["top", "bottom"]:
            raise Exception("Position must be 'top' or 'bottom'.")
        if alignment not in ["left", "center", "right"]:
            raise Exception("Alignment must be 'left', 'center', or 'right'.")
        if format not in ["page_x", "x"]:
            raise Exception("Format must be 'page_x' or 'x'.")
        
        for page_num in range(total_pages):
            page = doc[page_num]
            page_rect = page.rect
            font_size = 12
            # Choose text based on format
            text = f"Page {page_num + 1}" if format == "page_x" else f"{page_num + 1}"
            text_width = fitz.get_text_length(text, fontsize=font_size)
            margin = 20
            if position == "top":
                y = margin
            else:
                y = page_rect.height - margin - font_size
            if alignment == "left":
                x = margin
            elif alignment == "center":
                x = (page_rect.width - text_width) / 2
            else:
                x = page_rect.width - text_width - margin
            page.insert_text(
                (x, y),
                text,
                fontsize=font_size,
                color=(0, 0, 0),
                fontname="helv",
                overlay=True
            )
        output_buffer = io.BytesIO()
        doc.save(output_buffer)
        doc.close()
        output_buffer.seek(0)
        return output_buffer.getvalue()
    except Exception as e:
        logger.error(f"Page numbering error: {e}")
        return None
    finally:
        gc.collect()






def remove_background_rembg(image_bytes):
    """
    Remove background using rembg AI-based background remover.
    """
    try:
        output = remove(image_bytes)
        return io.BytesIO(output)
    except Exception as e:
        logger.error(f"Background removal failed: {str(e)}")
        raise ValueError(f"Background removal failed: {str(e)}")

def add_signature(pdf_bytes, signature_bytes, pages, size, position, alignment, remove_bg=False):
    try:
        logger.info("Starting add_signature function")

        if not pdf_bytes or not signature_bytes:
            raise ValueError("PDF or signature data is empty")
        if not pages:
            raise ValueError("No pages specified for signing")

        pdf_reader = PdfReader(io.BytesIO(pdf_bytes))
        num_pages = len(pdf_reader.pages)
        if num_pages == 0:
            raise ValueError("PDF has no pages")

        if not all(1 <= p <= num_pages for p in pages):
            raise ValueError(f"Page numbers out of range: {pages}")

        pdf_writer = PdfWriter()

        size_map = {
            'small': 100,
            'medium': 150,
            'large': 200
        }
        img_width = size_map.get(size, 150)

        # Process signature image based on remove_bg flag
        try:
            logger.info("Loading and processing signature image")
            if remove_bg:
                logger.info("Removing background from signature image using rembg")
                sig_io = remove_background_rembg(signature_bytes)
            else:
                sig_io = io.BytesIO(signature_bytes)
            img = ImageReader(sig_io)
            sig_width, sig_height = img.getSize()
            aspect = sig_height / float(sig_width) if sig_width else 1
            img_height = img_width * aspect
        except Exception as e:
            logger.error(f"Signature image processing failed: {str(e)}")
            raise ValueError(f"Invalid signature image: {str(e)}")

        for page_num in range(num_pages):
            page = pdf_reader.pages[page_num]
            page_width = float(page.mediabox.width)
            page_height = float(page.mediabox.height)

            if page_num + 1 in pages:
                packet = io.BytesIO()
                can = canvas.Canvas(packet, pagesize=(page_width, page_height))

                # Calculate position
                if position == 'top':
                    y = page_height - img_height - 50
                elif position == 'center':
                    y = (page_height - img_height) / 2
                else:  # bottom
                    y = 50

                if alignment == 'left':
                    x = 50
                elif alignment == 'center':
                    x = (page_width - img_width) / 2
                else:  # right
                    x = page_width - img_width - 50

                # Draw image with transparent background
                can.drawImage(img, x, y, width=img_width, height=img_height,
                              preserveAspectRatio=True, mask='auto')

                can.save()
                packet.seek(0)

                new_pdf = PdfReader(packet)
                if len(new_pdf.pages) == 0:
                    raise ValueError("Failed to create new PDF page")
                new_page = new_pdf.pages[0]

                try:
                    page.merge_page(new_page)
                except Exception as e:
                    logger.error(f"Page merge failed: {str(e)}")
                    raise RuntimeError(f"Page merge failed: {str(e)}")

            pdf_writer.add_page(page)

        output = io.BytesIO()
        pdf_writer.write(output)
        output.seek(0)
        result = output.read()
        if not result:
            raise ValueError("Generated PDF is empty")
        return result

    except Exception as e:
        logger.error(f"Error in signature processing: {str(e)}", exc_info=True)
        return None

def estimate_compression_sizes(pdf_bytes: bytes, custom_dpi: int, custom_quality: int) -> Optional[Dict[str, int]]:
    """Estimate output sizes for all compression presets."""
    try:
        sizes = {}
        presets = [
            ("high", 72, 20),
            ("medium", 100, 30),
            ("low", 120, 40),
            ("custom", custom_dpi, custom_quality)
        ]

        for preset_name, dpi, quality in presets:
            compressed_pdf = safe_compress_pdf(pdf_bytes, dpi, quality)
            if compressed_pdf is None:
                logger.error(f"Failed to compress for preset: {preset_name}")
                return None
            sizes[preset_name] = len(compressed_pdf)

        return sizes
    except Exception as e:
        logger.error(f"Size estimation error: {str(e)}")
        return None
    

# def estimate_compression_sizes(pdf_bytes: bytes, custom_dpi: int, custom_quality: int) -> Optional[Dict[str, int]]:
#     """Estimate output sizes for all compression presets."""
#     try:
#         sizes = {}
#         presets = [
#             ("high", 72, 20),
#             ("medium", 100, 30),
#             ("low", 120, 40),
#             ("custom", custom_dpi, custom_quality)
#         ]

#         for preset_name, dpi, quality in presets:
#             compressed_pdf = safe_compress_pdf(pdf_bytes, dpi, quality)
#             if compressed_pdf is None:
#                 logger.error(f"Failed to compress for preset: {preset_name}")
#                 return None
#             sizes[preset_name] = len(compressed_pdf)

#         return sizes
#     except Exception as e:
#         logger.error(f"Size estimation error: {str(e)}")
#         return None